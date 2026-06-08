import sys
import os
import json
import zipfile
import xml.etree.ElementTree as ET

# Force UTF-8 output so Cyrillic text is not garbled on Windows
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
else:
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import docx
except ImportError:
    docx = None


def extract_text_from_shape(shape, lines):
    """Recursively extract text from shapes, tables, and group shapes"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text and (not lines or lines[-1] != text):
                lines.append(text)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text and (not lines or lines[-1] != text):
                    lines.append(text)
    elif hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            extract_text_from_shape(sub_shape, lines)


def extract_slides_from_pptx(file_path):
    """Returns list of {num, image_path} for each slide using aspose.slides"""
    slides_info = []
    try:
        import aspose.slides as slides_module
    except ImportError:
        print("Error: aspose.slides not installed. Please pip install aspose-slides", file=sys.stderr)
        return []

    base_dir = os.path.dirname(file_path)
    base_name = os.path.splitext(os.path.basename(file_path))[0]

    try:
        pres = slides_module.Presentation(file_path)
        for i, slide in enumerate(pres.slides):
            out_filename = f"{base_name}_slide_{i+1}.png"
            out_path = os.path.join(base_dir, out_filename)
            # get_image(scaleX, scaleY) returns an IImage object
            # Saving with .png extension auto-detects format
            img = slide.get_image(2.0, 2.0)
            img.save(out_path)
            slides_info.append({
                'num': i + 1,
                'image_path': out_filename
            })
        return slides_info
    except Exception as e:
        print(f"Error extracting slides with aspose: {e}", file=sys.stderr)
        return []



def extract_paragraphs_from_docx(file_path):
    """Returns list of paragraph strings"""
    if docx is not None:
        try:
            paragraphs = []
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            text = para.text.strip()
                            if text:
                                paragraphs.append(text)
            return paragraphs
        except Exception as e:
            print(f"Warning: python-docx extraction failed: {e}", file=sys.stderr)

    # Fallback to manual XML parsing
    paragraphs = []
    try:
        docx_zip = zipfile.ZipFile(file_path)
        xml_content = docx_zip.read('word/document.xml')
        tree = ET.fromstring(xml_content)
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        for paragraph in tree.findall('.//w:p', ns):
            texts = [node.text for node in paragraph.findall('.//w:t', ns) if node.text]
            combined = ''.join(texts).strip()
            if combined:
                paragraphs.append(combined)
    except Exception as e:
        print(f"Error: Fallback docx XML parsing failed: {e}", file=sys.stderr)
    return paragraphs


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file"}))
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(json.dumps({"error": "not found"}))
        sys.exit(1)

    ext = file_path.lower().split('.')[-1]
    result = {"success": True, "type": ext}

    if ext == 'pptx':
        slides = extract_slides_from_pptx(file_path)
        result["slides"] = slides
        result["text"] = "Presentation images extracted"
    elif ext == 'docx':
        paragraphs = extract_paragraphs_from_docx(file_path)
        result["paragraphs"] = paragraphs
        result["text"] = '\n'.join(paragraphs)
    elif ext in ['jpg', 'jpeg', 'png', 'webp', 'gif']:
        if Image is None:
            result["optimized_path"] = os.path.basename(file_path)
        else:
            try:
                img = Image.open(file_path)
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                img.thumbnail((1920, 1920))
                base, _ = os.path.splitext(file_path)
                optimized_path = f"{base}_opt.jpg"
                img.save(optimized_path, "JPEG", quality=85)
                result["optimized_path"] = os.path.basename(optimized_path)
            except Exception as e:
                result["error"] = str(e)
                result["success"] = False
    else:
        result["error"] = "Unsupported file type"
        result["success"] = False

    print(json.dumps(result, ensure_ascii=False))
